"""
OKR × 業務計画 統合管理アプリ  ─  app.py （ローカル保存版）
=============================================================
S3の代わりにローカルフォルダ（data/）にJSONを保存します。
AWSアカウント不要。

インストール:
    pip install streamlit pandas plotly python-pptx kaleido

.streamlit/secrets.toml（最小構成）:
    [app]
    team_name = "プロダクトチーム"
    members   = ["田中 一郎", "鈴木 花子", "佐藤 健二", "山田 美咲"]
    admin_pin = "1234"
    data_dir  = "data"

複数人で共有する場合:
    data_dir に OneDrive や社内共有フォルダのパスを指定してください。
    例）data_dir = "C:/Users/Shared/okr_data"

起動:
    streamlit run app.py

データ保存先:
    data/
    ├── master_config.json       ← 四半期OKR
    └── plans/
        └── 2025-06/
            └── 2025-06_田中_一郎.json
"""

from __future__ import annotations

import io
import json
import datetime
import traceback
from pathlib import Path
from typing import Any

import gspread
import pandas as pd
import plotly.express as px
import streamlit as st
from google.oauth2.service_account import Credentials
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Cm

# ══════════════════════════════════════════════════════════════════════════════
# 定数
# ══════════════════════════════════════════════════════════════════════════════

PALETTE: list[dict] = [
    {"main": "#1B4F72", "light": "#D6EAF8", "text": "#0D2137"},
    {"main": "#1D6A45", "light": "#D5F5E3", "text": "#0B3A21"},
    {"main": "#6C3483", "light": "#EDE0F5", "text": "#4A235A"},
    {"main": "#7E5109", "light": "#FAE5D3", "text": "#4D320A"},
]
KR_COLORS   = ["#1B4F72", "#1D6A45", "#6C3483"]
KR_LABELS   = ["KR①", "KR②", "KR③"]
MAX_ACTIONS = 5
MAX_WALLS   = 3

DEFAULT_MASTER: dict = {
    "quarter":     "2025-Q2",
    "objective":   "",
    "locked":      False,
    "set_at":      "",
    "key_results": [
        {"id": "kr1", "label": "KR①", "text": "", "walls": []},
        {"id": "kr2", "label": "KR②", "text": "", "walls": []},
        {"id": "kr3", "label": "KR③", "text": "", "walls": []},
    ],
}

# ══════════════════════════════════════════════════════════════════════════════
# ページ設定 & CSS
# ══════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="OKR × 業務計画",
    page_icon="🌟",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
[data-testid="stSidebar"]{
    background:linear-gradient(175deg,#0D1B2A 0%,#1B4F72 70%,#1D6A45 100%);
}
[data-testid="stSidebar"] *{color:#C8E0FF !important;}
[data-testid="stSidebar"] .stSelectbox>div>div{
    background:rgba(200,224,255,.12)!important;
    border:1px solid rgba(200,224,255,.3)!important;
}
.northstar{background:#1B4F72;border-radius:14px;padding:1.1rem 1.6rem;
    margin-bottom:1.1rem;position:relative;overflow:hidden;}
.ns-orb1{position:absolute;right:-40px;top:-40px;width:150px;height:150px;
    border-radius:50%;background:rgba(243,156,18,.18);}
.ns-orb2{position:absolute;left:58%;bottom:-55px;width:120px;height:120px;
    border-radius:50%;background:rgba(29,106,69,.28);}
.ns-lock{position:absolute;top:11px;right:13px;border-radius:20px;padding:2px 11px;font-size:.7rem;}
.ns-lock.locked{background:rgba(243,156,18,.26);border:1px solid rgba(243,156,18,.5);color:#F9E08B;}
.ns-lock.unlocked{background:rgba(255,255,255,.14);border:1px solid rgba(255,255,255,.25);color:rgba(255,255,255,.7);}
.ns-eyebrow{font-size:.68rem;font-weight:700;letter-spacing:.12em;text-transform:uppercase;
    color:#F9E08B;margin-bottom:.35rem;position:relative;}
.ns-obj{font-size:1.1rem;font-weight:700;color:#fff;line-height:1.4;
    margin-bottom:.55rem;position:relative;}
.ns-obj.empty{color:rgba(255,255,255,.45);font-style:italic;font-weight:400;font-size:.9rem;}
.ns-krs{display:flex;gap:7px;flex-wrap:wrap;position:relative;}
.ns-kr{background:rgba(255,255,255,.13);border:1px solid rgba(255,255,255,.22);
    border-radius:20px;padding:3px 12px;font-size:.74rem;color:#C8E0FF;}
.g-info{background:var(--color-background-info);border-left:3px solid var(--color-border-info);
    color:var(--color-text-info);border-radius:0 8px 8px 0;
    padding:.6rem .9rem;font-size:.8rem;line-height:1.55;margin:.4rem 0;}
.g-ok{background:#EAF3DE;border-left:3px solid #7dcea0;color:#1A5632;
    border-radius:0 8px 8px 0;padding:.55rem .9rem;font-size:.78rem;margin:.3rem 0;}
.g-ng{background:#FDEDEC;border-left:3px solid #f1948a;color:#7B241C;
    border-radius:0 8px 8px 0;padding:.55rem .9rem;font-size:.78rem;margin:.3rem 0;}
.g-warn{background:#FEF9E7;border-left:3px solid #F39C12;color:#7D6608;
    border-radius:0 8px 8px 0;padding:.6rem .9rem;font-size:.78rem;margin:.4rem 0;}
.prog-wrap{display:flex;align-items:center;margin-bottom:1.1rem;}
.prog-step{display:flex;align-items:center;gap:6px;padding:.42rem .9rem;
    border-radius:20px;font-size:.78rem;font-weight:500;white-space:nowrap;}
.prog-step.done{background:#EAF3DE;color:#1A5632;}
.prog-step.active{background:#1B4F72;color:#fff;}
.prog-step.todo{background:var(--color-background-secondary);color:var(--color-text-secondary);}
.prog-line{flex:1;height:2px;background:var(--color-border-tertiary);min-width:14px;}
.prog-line.filled{background:#1A5632;}
.ltree{background:var(--color-background-secondary);border:0.5px solid var(--color-border-tertiary);
    border-radius:var(--border-radius-lg);padding:1rem 1.15rem;}
.ltree-hdr{font-size:.68rem;font-weight:600;letter-spacing:.09em;text-transform:uppercase;
    color:var(--color-text-secondary);margin-bottom:.8rem;display:flex;align-items:center;gap:6px;}
.lt-row{display:flex;gap:8px;align-items:flex-start;margin-bottom:.5rem;}
.lt-ico{width:22px;height:22px;border-radius:5px;display:flex;align-items:center;
    justify-content:center;font-size:.65rem;font-weight:700;color:#fff;flex-shrink:0;margin-top:1px;}
.lt-body{font-size:.8rem;color:var(--color-text-primary);line-height:1.45;flex:1;}
.lt-sub{font-size:.7rem;color:var(--color-text-secondary);margin-top:2px;}
.lt-indent{margin-left:15px;padding-left:13px;border-left:2px solid var(--color-border-tertiary);}
.lt-empty{font-size:.75rem;color:var(--color-text-tertiary);font-style:italic;}
.check-list{display:flex;flex-direction:column;gap:5px;margin:.5rem 0;}
.check-item{display:flex;align-items:flex-start;gap:7px;font-size:.82rem;}
.member-grid{display:grid;grid-template-columns:repeat(4,1fr);gap:10px;margin-bottom:1rem;}
.member-card{border-radius:10px;padding:.8rem;text-align:center;border:1.5px solid;}
.mc-av{width:36px;height:36px;border-radius:50%;display:flex;align-items:center;
    justify-content:center;font-size:.85rem;font-weight:600;margin:0 auto .4rem;}
.mc-name{font-size:.82rem;font-weight:600;margin-bottom:3px;}
.mc-badge{font-size:.7rem;padding:2px 8px;border-radius:20px;display:inline-block;font-weight:500;}
.stat-row{display:flex;gap:10px;margin-bottom:1.1rem;flex-wrap:wrap;}
.stat{background:var(--color-background-secondary);border-radius:var(--border-radius-md);
    padding:.7rem 1rem;text-align:center;flex:1;min-width:90px;}
.stat-n{font-size:1.55rem;font-weight:600;color:#1B4F72;}
.stat-l{font-size:.72rem;color:var(--color-text-secondary);margin-top:2px;}
.dash-card{background:var(--color-background-primary);border:0.5px solid var(--color-border-tertiary);
    border-radius:var(--border-radius-lg);padding:1rem 1.2rem;margin-bottom:.85rem;border-left:4px solid;}
.dash-hdr{display:flex;align-items:center;gap:10px;margin-bottom:.65rem;}
.dash-av{width:34px;height:34px;border-radius:50%;display:flex;align-items:center;
    justify-content:center;font-size:.8rem;font-weight:600;flex-shrink:0;}
.kr-block{background:var(--color-background-secondary);border-radius:var(--border-radius-md);
    padding:.65rem .85rem;margin-bottom:.5rem;}
.action-list-item{display:flex;gap:7px;align-items:flex-start;padding:.28rem 0;
    border-bottom:0.5px solid var(--color-border-tertiary);font-size:.78rem;}
.action-list-item:last-child{border-bottom:none;}
.a-num{width:18px;height:18px;border-radius:50%;display:flex;align-items:center;
    justify-content:center;font-size:.62rem;font-weight:600;color:#fff;flex-shrink:0;margin-top:1px;}
.fb-panel{background:#FEF9E7;border:0.5px solid #F9E79F;border-radius:var(--border-radius-lg);
    padding:1rem 1.2rem;margin-bottom:1rem;}
.fb-panel h4{font-size:.88rem;font-weight:600;color:#7D6608;margin-bottom:.55rem;}
.fb-item{display:flex;gap:8px;padding:.3rem 0;border-bottom:0.5px solid rgba(249,231,159,.7);
    font-size:.8rem;color:#6E4D0B;align-items:flex-start;}
.fb-item:last-child{border-bottom:none;}
.fb-ico{font-size:14px;flex-shrink:0;margin-top:1px;}
.flow-row{display:flex;align-items:stretch;gap:0;margin-bottom:.9rem;flex-wrap:nowrap;}
.flow-box{flex:1;border-radius:10px;padding:.85rem .9rem;text-align:center;min-width:0;}
.flow-arr{display:flex;align-items:center;padding:0 5px;
    color:var(--color-text-secondary);font-size:1rem;flex-shrink:0;}
.local-badge{background:#E8F4F8;border:0.5px solid #85B7EB;border-radius:7px;
    padding:5px 11px;font-size:.8rem;color:#0C447C;margin-bottom:.6rem;}
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# 設定読み込み
# ══════════════════════════════════════════════════════════════════════════════

@st.cache_data(ttl=0)
def load_config() -> dict:
    try:
        members   = list(st.secrets["app"]["members"])
        team_name = str(st.secrets["app"].get("team_name", "チーム"))
        admin_pin = str(st.secrets["app"].get("admin_pin", "1234"))
        sheet_id  = str(st.secrets["app"].get("spreadsheet_id", ""))
    except Exception:
        members   = ["田中 一郎", "鈴木 花子", "佐藤 健二", "山田 美咲"]
        team_name = "プロダクトチーム"
        admin_pin = "1234"
        sheet_id  = ""
    return dict(members=members, team_name=team_name,
                admin_pin=admin_pin, sheet_id=sheet_id)


CFG     = load_config()
MEMBERS = CFG["members"]

# ローカルフォールバック用（Sheets未設定時）
BASE_DIR = Path(__file__).parent / "data"


# ══════════════════════════════════════════════════════════════════════════════
# Google Sheets クライアント
# ══════════════════════════════════════════════════════════════════════════════

@st.cache_resource(ttl=600)
def get_gsheet_client():
    """gspreadクライアントを返す。未設定の場合はNone。"""
    try:
        creds_dict = dict(st.secrets["gcp"])
        if "private_key" in creds_dict:
            creds_dict["private_key"] = creds_dict["private_key"].replace("\\n", "\n")
        scopes = [
            "https://spreadsheets.google.com/feeds",
            "https://www.googleapis.com/auth/drive",
        ]
        creds  = Credentials.from_service_account_info(creds_dict, scopes=scopes)
        client = gspread.authorize(creds)
        client.open_by_key(dict(st.secrets["app"]).get("spreadsheet_id",""))
        return client
    except Exception as e:
        return f"ERROR:{e}"


def _get_or_create_sheet(client, title: str):
    """スプレッドシート内のシートを取得。なければ作成。"""
    ss = client.open_by_key(CFG["sheet_id"])
    try:
        return ss.worksheet(title)
    except gspread.WorksheetNotFound:
        ws = ss.add_worksheet(title=title, rows=10, cols=2)
        ws.append_row(["key", "value"])
        return ws


@st.cache_data(ttl=120)
def _cached_sheet_records(sheet_id: str, sheet_title: str) -> list[dict]:
    """シートの全レコードをキャッシュ付きで取得（2分間キャッシュ）。"""
    try:
        client = _valid_client()
        if not client:
            return []
        ws = _get_or_create_sheet(client, sheet_title)
        return ws.get_all_records()
    except Exception:
        return []


def _invalidate_sheet_cache(sheet_id: str, sheet_title: str):
    """保存後にキャッシュを無効化する。"""
    _cached_sheet_records.clear()


def _sheets_get(client, sheet_title: str, key: str) -> dict | None:
    """シートからkeyに対応するJSONを取得。"""
    try:
        rows = _cached_sheet_records(CFG["sheet_id"], sheet_title)
        for row in rows:
            if str(row.get("key","")) == key:
                return json.loads(row["value"])
        return None
    except Exception:
        return None


def _sheets_set(client, sheet_title: str, key: str, data: dict) -> bool:
    """シートにkeyとJSONを保存（既存行は上書き）。"""
    try:
        ws    = _get_or_create_sheet(client, sheet_title)
        rows  = ws.get_all_records()
        value = json.dumps(data, ensure_ascii=False)
        for i, row in enumerate(rows, start=2):
            if str(row.get("key","")) == key:
                ws.update_cell(i, 2, value)
                _invalidate_sheet_cache(CFG["sheet_id"], sheet_title)
                return True
        ws.append_row([key, value])
        _invalidate_sheet_cache(CFG["sheet_id"], sheet_title)
        return True
    except Exception as e:
        st.error(f"Sheets保存エラー: {e}")
        return False


def _sheets_delete(client, sheet_title: str, key: str) -> bool:
    """シートからkeyの行を削除。"""
    try:
        ws   = _get_or_create_sheet(client, sheet_title)
        rows = ws.get_all_records()
        for i, row in enumerate(rows, start=2):
            if str(row.get("key","")) == key:
                ws.delete_rows(i)
                _invalidate_sheet_cache(CFG["sheet_id"], sheet_title)
                return True
        return True
    except Exception:
        return False


def _sheets_list(client, sheet_title: str, prefix: str) -> list[dict]:
    """指定prefixで始まるkeyのJSON一覧を返す。"""
    try:
        rows = _cached_sheet_records(CFG["sheet_id"], sheet_title)
        result = []
        for row in rows:
            k = str(row.get("key",""))
            if k.startswith(prefix):
                try:
                    result.append(json.loads(row["value"]))
                except Exception:
                    pass
        return result
    except Exception:
        return []


# ── ローカルフォールバック ────────────────────────────────────────────────

def _local_get(key: str) -> dict | None:
    p = BASE_DIR / f"{key.replace('/', '_')}.json"
    try:
        return json.loads(p.read_text(encoding="utf-8")) if p.exists() else None
    except Exception:
        return None


def _local_set(key: str, data: dict) -> bool:
    p = BASE_DIR / f"{key.replace('/', '_')}.json"
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        return True
    except Exception as e:
        st.error(f"ローカル保存エラー: {e}")
        return False


def _local_delete(key: str) -> bool:
    p = BASE_DIR / f"{key.replace('/', '_')}.json"
    if p.exists():
        p.unlink()
    return True


def _local_list(prefix: str) -> list[dict]:
    safe = prefix.replace("/", "_")
    result = []
    for p in sorted(BASE_DIR.glob(f"{safe}*.json")):
        try:
            result.append(json.loads(p.read_text(encoding="utf-8")))
        except Exception:
            pass
    return result


# ── 統合I/O（Sheets優先、フォールバックはローカル） ────────────────────────

MASTER_KEY   = "master_config"
PLANS_SHEET  = "plans"


def plan_key(month_str: str, member: str) -> str:
    safe = member.replace(" ", "_").replace("　", "_")
    return f"{month_str}/{safe}"


def _valid_client():
    """有効なgspreadクライアントを返す。エラーまたは未設定の場合はNone。"""
    c = get_gsheet_client()
    return c if c and not isinstance(c, str) else None


def io_get_master() -> dict:
    client = _valid_client()
    if client and CFG["sheet_id"]:
        data = _sheets_get(client, "master", MASTER_KEY)
        return data if data else DEFAULT_MASTER.copy()
    return _local_get(MASTER_KEY) or DEFAULT_MASTER.copy()


def io_save_master(data: dict) -> bool:
    client = _valid_client()
    if client and CFG["sheet_id"]:
        return _sheets_set(client, "master", MASTER_KEY, data)
    return _local_set(MASTER_KEY, data)


def io_get_plan(month_str: str) -> dict | None:
    """チームの月次プランを取得。"""
    client = _valid_client()
    key    = f"{month_str}/team"
    if client and CFG["sheet_id"]:
        return _sheets_get(client, PLANS_SHEET, key)
    return _local_get(key)


def io_save_plan(month_str: str, data: dict) -> bool:
    """チームの月次プランを保存。"""
    client = _valid_client()
    key    = f"{month_str}/team"
    if client and CFG["sheet_id"]:
        return _sheets_set(client, PLANS_SHEET, key, data)
    return _local_set(key, data)


def io_delete_plan(month_str: str) -> bool:
    """チームの月次プランを削除。"""
    client = _valid_client()
    key    = f"{month_str}/team"
    if client and CFG["sheet_id"]:
        return _sheets_delete(client, PLANS_SHEET, key)
    return _local_delete(key)


def io_list_plans(month_str: str) -> list[dict]:
    """チームプランを1件のリストで返す（互換性維持）。"""
    data = io_get_plan(month_str)
    return [data] if data else []


PRIORITY_SHEET = "priorities"


def io_get_priorities(month_str: str) -> dict:
    client = _valid_client()
    if client and CFG["sheet_id"]:
        data = _sheets_get(client, PRIORITY_SHEET, month_str)
        return data if data else {}
    return _local_get(f"priorities_{month_str}") or {}


def io_save_priorities(month_str: str, data: dict) -> bool:
    client = _valid_client()
    if client and CFG["sheet_id"]:
        return _sheets_set(client, PRIORITY_SHEET, month_str, data)
    return _local_set(f"priorities_{month_str}", data)


# ══════════════════════════════════════════════════════════════════════════════
# ヘルパー
# ══════════════════════════════════════════════════════════════════════════════

def mpal(name: str) -> dict:
    idx = MEMBERS.index(name) if name in MEMBERS else 0
    return PALETTE[idx % len(PALETTE)]


def initials(name: str) -> str:
    parts = name.split()
    return parts[0][0] + (parts[1][0] if len(parts) > 1 else "")


def blank_action() -> dict:
    today = datetime.date.today()
    return {"text": "", "start": today.isoformat(),
            "end": (today + datetime.timedelta(days=30)).isoformat()}


def blank_wall_actions(walls: list[str]) -> list[dict]:
    """壁リストに対応した空のアクション辞書を返す。
    [{wall_text, actions:[{text,start,end}]}]
    """
    return [{"wall_text": w, "actions": [blank_action()]} for w in walls]


def _init_session():
    today = datetime.date.today()
    for k, v in dict(
        month_str    = today.strftime("%Y-%m"),
        plan_step    = 0,
        plan_kr_idx  = 0,
        plan_actions = [],
        admin_auth   = False,
        team_data    = None,
    ).items():
        if k not in st.session_state:
            st.session_state[k] = v


# ══════════════════════════════════════════════════════════════════════════════
# 北極星バナー
# ══════════════════════════════════════════════════════════════════════════════

def render_north_star(master: dict):
    locked   = master.get("locked", False)
    obj      = master.get("objective", "").strip()
    krs      = master.get("key_results", [])
    qtr      = master.get("quarter", "")
    lock_cls = "locked" if locked else "unlocked"
    lock_lbl = f"🔒 {qtr} 確定済み" if locked else f"📝 {qtr} 未確定"
    kr_html  = " ".join(
        f'<span class="ns-kr">{kr["label"]}：{kr["text"]}</span>'
        for kr in krs if kr.get("text")
    ) if obj else ""
    obj_html = (
        f'<div class="ns-obj">{obj}</div>' if obj
        else '<div class="ns-obj empty">まだ設定されていません ─ STRATEGYタブで入力してください</div>'
    )
    st.markdown(f"""
<div class="northstar">
  <div class="ns-orb1"></div><div class="ns-orb2"></div>
  <span class="ns-lock {lock_cls}">{lock_lbl}</span>
  <div class="ns-eyebrow">🌟 今期の北極星 Objective</div>
  {obj_html}
  <div class="ns-krs">{kr_html}</div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# プログレスバー
# ══════════════════════════════════════════════════════════════════════════════

def render_progress(step: int):
    steps = ["🎯 KRを選ぶ", "🔍 壁とアクションを決める"]
    parts = []
    for i, label in enumerate(steps):
        cls = "done" if i < step else ("active" if i == step else "todo")
        parts.append(f'<div class="prog-step {cls}">{label}</div>')
        if i < len(steps) - 1:
            parts.append(f'<div class="prog-line {"filled" if i < step else ""}"></div>')
    st.markdown(f'<div class="prog-wrap">{"".join(parts)}</div>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# ロジックツリー
# ══════════════════════════════════════════════════════════════════════════════

def render_logic_tree(master: dict, kr_idx: int,
                      wall_actions: list[dict], pal: dict):
    """wall_actions = [{wall_text, actions:[{text,start,end}]}]"""
    krs = master.get("key_results", [])
    kr  = krs[kr_idx] if kr_idx < len(krs) else {}

    if wall_actions:
        walls_html = ""
        for ii, wa in enumerate(wall_actions):
            filled_actions = [a for a in wa.get("actions",[]) if a.get("text","").strip()]
            actions_html = "".join(
                f'<div class="lt-indent" style="margin-top:.3rem;">'
                f'<div class="lt-row">'
                f'<div class="lt-ico" style="background:{pal["main"]};">A{ia+1}</div>'
                f'<div class="lt-body"><span style="font-weight:600;">{a["text"]}</span>'
                f'<div class="lt-sub">{a.get("start","")} → {a.get("end","")}</div>'
                f'</div></div></div>'
                for ia, a in enumerate(filled_actions)
            ) if filled_actions else (
                f'<div class="lt-indent" style="margin-top:.3rem;">'
                f'<div class="lt-row"><div class="lt-ico" style="background:{pal["main"]};">A</div>'
                f'<div class="lt-body lt-empty">アクションを入力してください</div>'
                f'</div></div>'
            )
            walls_html += (
                f'<div class="lt-row" style="margin-top:.35rem;">'
                f'<div class="lt-ico" style="background:#F39C12;">壁{ii+1}</div>'
                f'<div class="lt-body" style="font-weight:600;">{wa["wall_text"]}</div>'
                f'</div>{actions_html}'
            )
    else:
        walls_html = (
            '<div class="lt-row" style="margin-top:.35rem;">'
            '<div class="lt-ico" style="background:#F39C12;">壁</div>'
            '<div class="lt-body lt-empty">KRを選ぶと壁が表示されます</div>'
            '</div>'
        )

    st.markdown(f"""
<div class="ltree">
  <div class="ltree-hdr">🌲 ロジックツリー（リアルタイム）</div>
  <div class="lt-row">
    <div class="lt-ico" style="background:#1B4F72;">KR</div>
    <div class="lt-body" style="font-weight:700;">{kr.get("label","KR")}：{kr.get("text","")}</div>
  </div>
  <div class="lt-indent">{walls_html}</div>
</div>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# HOME タブ
# ══════════════════════════════════════════════════════════════════════════════

def render_home():
    st.markdown("""
<div class="flow-row">
  <div class="flow-box" style="background:#EBF5FB;border:1.5px solid #1B4F72;">
    <div style="font-size:1.35rem;margin-bottom:.25rem">🌟</div>
    <div style="font-size:.7rem;font-weight:700;letter-spacing:.06em;text-transform:uppercase;color:#1B4F72">Objective</div>
    <div style="font-size:.78rem;margin-top:.2rem;line-height:1.4;color:var(--color-text-secondary)">ワクワクする<br>定性的な<b>目的地</b></div>
  </div>
  <div class="flow-arr">→</div>
  <div class="flow-box" style="background:#EAF4FB;border:1.5px solid #2E86C1;">
    <div style="font-size:1.35rem;margin-bottom:.25rem">📏</div>
    <div style="font-size:.7rem;font-weight:700;letter-spacing:.06em;text-transform:uppercase;color:#2E86C1">Key Results</div>
    <div style="font-size:.78rem;margin-top:.2rem;line-height:1.4;color:var(--color-text-secondary)">達成を測る<br>定量的な<b>指標</b>（3つ）</div>
  </div>
  <div class="flow-arr">→</div>
  <div class="flow-box" style="background:#FEF9E7;border:1.5px solid #F39C12;">
    <div style="font-size:1.35rem;margin-bottom:.25rem">🔍</div>
    <div style="font-size:.7rem;font-weight:700;letter-spacing:.06em;text-transform:uppercase;color:#B7770D">壁（課題）</div>
    <div style="font-size:.78rem;margin-top:.2rem;line-height:1.4;color:var(--color-text-secondary)">KRに届かない<br><b>真の原因</b></div>
  </div>
  <div class="flow-arr">→</div>
  <div class="flow-box" style="background:#EAFAF1;border:1.5px solid #1D6A45;">
    <div style="font-size:1.35rem;margin-bottom:.25rem">⚡</div>
    <div style="font-size:.7rem;font-weight:700;letter-spacing:.06em;text-transform:uppercase;color:#1D6A45">Action（行動）</div>
    <div style="font-size:.78rem;margin-top:.2rem;line-height:1.4;color:var(--color-text-secondary)">今月<b>やり切る</b><br>具体的な一手（複数可）</div>
  </div>
  <div class="flow-arr">→</div>
  <div class="flow-box" style="background:#F4ECF7;border:1.5px solid #7D3C98;">
    <div style="font-size:1.35rem;margin-bottom:.25rem">📋</div>
    <div style="font-size:.7rem;font-weight:700;letter-spacing:.06em;text-transform:uppercase;color:#7D3C98">Backlog起票</div>
    <div style="font-size:.78rem;margin-top:.2rem;line-height:1.4;color:var(--color-text-secondary)">保存後に<br><b>ワンクリック起票</b></div>
  </div>
</div>
""", unsafe_allow_html=True)
    st.markdown('<div class="g-info"><b>なぜこの順番？</b>　KR → 壁（課題） → アクション の順に強制的に考えさせる設計。アクションを保存するとBacklogへのワンクリック起票リンクが表示されます。</div>', unsafe_allow_html=True)

    # ── 操作ガイドボタン ──────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("### 📖 操作ガイド")
    gc1, gc2 = st.columns(2)

    with gc1:
        if st.button("👤 メンバー向けガイドを見る", use_container_width=True):
            st.session_state["show_member_guide"] = True

    with gc2:
        if st.button("👔 マネジャー向けガイドを見る", use_container_width=True):
            st.session_state["show_manager_guide"] = True

    # メンバー向けガイド ダイアログ
    if st.session_state.get("show_member_guide"):
        @st.dialog("👤 メンバー向け 操作ガイド", width="large")
        def member_guide_dialog():
            st.markdown("### 1. アクセス方法")
            st.markdown("ブラウザでアプリのURLを開くだけです。インストール不要。左のサイドバーで **自分の名前** を必ず選んでください。")
            st.warning("⚠️ 名前を間違えると他の人のデータを上書きしてしまいます。")

            st.markdown("---")
            st.markdown("### 2. 月次計画の入力（PLANタブ）")
            st.markdown("「📝 PLAN」タブを開き、2ステップで入力します。")

            st.markdown("#### STEP 1：KRを選ぶ")
            st.markdown("今月最も重要なKRを1つ選んで「選ぶ」ボタンを押してください。")

            st.markdown("#### STEP 2：壁とアクションを決める")
            st.markdown("**壁（課題）を入力する**\n「なぜこのKRに今届いていないのか？」の原因を1文で書いてください。")
            st.markdown('<div class="g-ok">✅ 「インタビューがゼロのため施策が推測止まりになっている」</div>', unsafe_allow_html=True)
            st.markdown('<div class="g-ng">❌ 「NPSが低い」← 現象ではなく原因を書く</div>', unsafe_allow_html=True)

            st.markdown("**アクションを入力する**\n各壁の下に「いつまでに・何を・どれだけ」を含む行動を入力してください。")
            st.markdown('<div class="g-ok">✅ 「6/15までにインタビュー10件実施しSlackで共有する」</div>', unsafe_allow_html=True)
            st.markdown('<div class="g-ng">❌ 「顧客の声を聞く」← いつ？何件？</div>', unsafe_allow_html=True)

            st.markdown("""
**壁・アクションの追加・削除**
- 「＋ 壁（課題）を追加する」で壁を最大3件まで追加できます
- 「＋ アクションを追加」で各壁に最大5件のアクションを追加できます
- 🗑 ボタンで不要な壁・アクションを削除できます
""")
            st.markdown("入力が終わったら「💾 保存する」を押してください。保存後にBacklog起票画面が表示されます。")

            st.markdown("---")
            st.markdown("### 3. 複数のKRに入力したい場合")
            st.markdown("「別のKRも入力する」ボタンでKR選択画面に戻ります。入力済みのデータは保持されます。")

            st.markdown("---")
            st.markdown("### 4. よくある質問")
            st.markdown("""
**Q. 保存前にブラウザを閉じてしまったら？**
入力内容は消えます。こまめに「💾 保存する」を押しましょう。

**Q. 入力を修正したい**
再度PLANタブを開くと前回のデータが読み込まれます。修正して保存すると上書きされます。

**Q. 入力期限は？**
毎月1〜3日を目安に入力してください。
""")
            st.markdown("---")
            if st.button("閉じる", use_container_width=True):
                st.session_state["show_member_guide"] = False
                st.rerun()

        member_guide_dialog()

    # マネジャー向けガイド ダイアログ
    if st.session_state.get("show_manager_guide"):
        @st.dialog("👔 マネジャー向け 操作ガイド", width="large")
        def manager_guide_dialog():
            st.markdown("### 1. 月次運用カレンダー")
            st.markdown("""
| タイミング | 作業 | 使う機能 |
|-----------|------|---------|
| 四半期初め | OKRを確定する | STRATEGYタブ |
| 月初 1〜3日 | メンバーの入力を促す | （メンバーがPLANタブで入力） |
| 月初 MTG当日 | 全員の計画を確認 | DASHBOARDタブ |
| 月初 MTG当日 | PPTX資料を生成 | DASHBOARDタブ → PPTX出力 |
| 月末 | 振り返りを実施 | DASHBOARDタブ |
""")
            st.markdown("---")
            st.markdown("### 2. 四半期OKRの設定（STRATEGYタブ）")
            st.markdown("""
1. 「🏛️ STRATEGY」タブを開き、管理者PINを入力して認証
2. **Objective**：チームがワクワクする定性的な目標を入力（右側のチェックリストが全✅になるのが理想）
3. **Key Results**：月末に○か✕か判定できる数値指標を3つ入力（70〜80%達成が理想的な難易度）
4. 「🔒 OKRを確定保存する」を押すとロック → 北極星バナーに反映
""")
            st.warning("⚠️ 確定後は四半期中は編集不可。変更は「ロック解除」ボタン＋チーム全員の合意が必要です。")

            st.markdown("---")
            st.markdown("### 3. ダッシュボードの活用（DASHBOARDタブ）")
            st.markdown("""
- 「🔄 データを読み込む」で全員の最新データを取得
- 提出状況カードで未提出メンバーを確認
- 「マネジャー向け レビューの着眼点」を展開してフィードバック観点を確認
""")
            st.markdown("**フィードバックの5観点**")
            st.markdown("""
| 観点 | チェックポイント |
|-----|----------------|
| 🔥 野心度 | KRは70〜80%達成が理想の難易度か？ |
| 🔗 ロジック | 壁→アクションの因果関係は通っているか？ |
| ⚡ リソース | 全アクション合計で1ヶ月に現実的な量か？ |
| 🤝 連携 | 他メンバーとの重複・依存はないか？ |
| 🎯 整合性 | 全員がObjectiveを向いているか？ |
""")
            st.markdown("---")
            st.markdown("### 4. PPTX資料の生成")
            st.markdown("""
1. DASHBOARDタブで「🔄 データを読み込む」
2. 「🚀 PPTXを生成する」を押す
3. 数秒後に「⬇️ PPTXをダウンロード」が表示される

**構成：** 表紙 → チームOKRサマリー → メンバー別詳細 → 統合ガントチャート
""")
            st.markdown("---")
            if st.button("閉じる", use_container_width=True):
                st.session_state["show_manager_guide"] = False
                st.rerun()

        manager_guide_dialog()

    st.markdown("---")
    st.markdown("### 月次の運用サイクル")
    cols = st.columns(4)
    for col, (n, title, desc) in zip(cols, [
        (1, "四半期初め",  "管理者がSTRATEGYでOKRを確定。"),
        (2, "月初 1〜3日", "各メンバーがPLANタブで計画を入力・保存。"),
        (3, "月初 MTG",   "DASHBOARDで全員の計画を確認・PPTX生成。"),
        (4, "月末",        "振り返りを行い、翌月の計画に活かす。"),
    ]):
        with col:
            st.markdown(
                f'<div style="background:var(--color-background-secondary);border-radius:10px;padding:.8rem;text-align:center;">'
                f'<div style="font-size:1.5rem;font-weight:600;color:#1B4F72;">{n}</div>'
                f'<div style="font-size:.82rem;font-weight:600;color:var(--color-text-primary);">{title}</div>'
                f'<div style="font-size:.75rem;color:var(--color-text-secondary);margin-top:3px;line-height:1.45;">{desc}</div></div>',
                unsafe_allow_html=True,
            )
    st.markdown("---")
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown("**🌟 良いObjectiveとは**")
        st.markdown('<div class="g-ok">✅ 「顧客が思わず人に紹介したくなるサービスにする」</div>', unsafe_allow_html=True)
        st.markdown('<div class="g-ng">❌ 「売上を上げる」← 数値はKRへ</div>', unsafe_allow_html=True)
    with c2:
        st.markdown("**📏 良いKey Resultとは**")
        st.markdown('<div class="g-ok">✅ 「月次NPS 40以上」「解約率 1%以下」</div>', unsafe_allow_html=True)
        st.markdown('<div class="g-ng">❌ 「顧客満足度を高める」← 測れない</div>', unsafe_allow_html=True)
    with c3:
        st.markdown("**⚡ 良いActionとは**")
        st.markdown('<div class="g-ok">✅ 「6/15までにインタビュー10件実施」</div>', unsafe_allow_html=True)
        st.markdown('<div class="g-ng">❌ 「顧客の声を聞く」← いつ？何件？</div>', unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# STRATEGY タブ
# ══════════════════════════════════════════════════════════════════════════════

def render_strategy(master: dict):
    if not st.session_state.admin_auth:
        st.markdown("### 🔐 管理者モードの認証")
        st.markdown('<div class="g-info"><b>このタブはチームの四半期OKRを確定する専用画面です。</b><br>全員で画面を共有しながら入力し、合意の上で「確定保存」してください。</div>', unsafe_allow_html=True)
        with st.form("admin_login_form"):
            pin = st.text_input("管理者PINコードを入力", type="password", key="admin_pin_input")
            if st.form_submit_button("認証する", use_container_width=True):
                if pin == CFG["admin_pin"] or not pin:
                    st.session_state.admin_auth = True
                    st.rerun()
                else:
                    st.error("PINコードが違います。")
        return

    locked = master.get("locked", False)
    if locked:
        st.success(f"✅ このOKRは **{master.get('quarter','')}** の期間中、確定済みです。")
        st.markdown('<div class="g-warn">🔒 確定済みのOKRは四半期中は編集できません。変更が必要な場合はチーム全員で合意し、下のボタンでロックを解除してください。</div>', unsafe_allow_html=True)
        with st.expander("確定済みOKRの内容を確認"):
            st.markdown(f"**Objective：** {master.get('objective','')}")
            for kr in master.get("key_results", []):
                st.markdown(f"- **{kr['label']}：** {kr.get('text','')}")
        if st.button("🔓 ロックを解除して編集する（要チーム合意）", type="secondary"):
            master["locked"] = False
            io_save_master(master)
            st.session_state.cached_master = master
            st.toast("ロックを解除しました。", icon="🔓")
            st.rerun()
        return

    st.markdown('<div class="g-info">💡 チーム全員で画面を見ながら話し合いましょう。「ワクワクするか？」「数値で測れるか？」を声に出しながら入力してください。</div>', unsafe_allow_html=True)

    with st.container(border=True):
        st.markdown("**対象四半期**")
        quarter = st.text_input("四半期", value=master.get("quarter","2025-Q2"), label_visibility="collapsed")

    with st.container(border=True):
        st.markdown("**🌟 Objective（チームの目的地）**　─　定性的・感情的な言葉で")
        col_in, col_check = st.columns([3, 2])
        with col_in:
            objective = st.text_area("Objective",
                value=master.get("objective",""), height=100,
                placeholder="例）顧客が思わず人に紹介したくなるプロダクトを作る",
                help="数値はNG！「〜する」「〜になる」という完成形で。",
                label_visibility="collapsed")
        with col_check:
            obj_v = objective.strip()
            checks = [
                ("読んだ時にワクワクする",       len(obj_v) > 8),
                ("数字が含まれていない",          not any(c.isdigit() for c in obj_v)),
                ("「〜する/になる」の完成形",     obj_v.endswith(("する","なる","きる","げる","れる","つ","む","ぐ"))),
                ("15〜70文字程度",               15 <= len(obj_v) <= 70),
            ]
            st.markdown('<div class="check-list">', unsafe_allow_html=True)
            for label, ok in checks:
                color = "#196F3D" if ok else "var(--color-text-secondary)"
                st.markdown(f'<div class="check-item" style="color:{color};">{"✅" if ok else "⬜"} {label}</div>', unsafe_allow_html=True)
            st.markdown("</div>", unsafe_allow_html=True)

    with st.container(border=True):
        st.markdown("**📏 Key Results と 壁（課題）**　─　KRごとにチームで合意した壁を設定してください")
        st.markdown('<div class="g-info">KRは「Objectiveが達成された証拠」。壁は「チームで議論したKRに届かない理由」です。各メンバーはこの壁に対するアクションを考えます。</div>', unsafe_allow_html=True)

        # 壁の考え方ヒント
        if st.button("💡 壁（課題）の考え方ヒント", use_container_width=False):
            st.session_state["show_wall_hint"] = True
        if st.session_state.get("show_wall_hint"):
            @st.dialog("💡 壁（課題）の考え方ヒント")
            def wall_hint_dialog():
                st.markdown("### 壁とは何か？")
                st.markdown("「なぜ今月このKRに届いていないのか？」の**根本原因**です。表面的な現象ではなく、その奥にある原因を掘り下げてください。")
                st.markdown("---")
                st.markdown("### ✅ 良い壁の書き方")
                st.markdown("""
**「〇〇が不足しているため、〜できていない」** の形で書くと論理が明確になります。

- 「顧客インタビューがゼロのため、ペインが推測止まりになっており施策が的外れになっている」
- 「提案資料に数字の根拠がなく、顧客が意思決定に踏み切れていない」
- 「チュートリアルが長すぎて、新規ユーザーが離脱している」
""")
                st.markdown("### ❌ 避けてほしい書き方")
                st.markdown("""
- 「NPSが低い」← 現象を書いても意味がない。**なぜ**低いのかを書く
- 「頑張る」← 壁ではなく意気込みになっている
- 「時間がない」← 本当の原因をもう一段掘り下げる
""")
                st.markdown("---")
                st.markdown("### 🔍 考えるための問いかけ")
                st.markdown("""
1. **「なぜKRに届いていないのか？」** を5回繰り返してみる
2. **「もし明日このKRを達成するとしたら、何が邪魔をしているか？」**
3. **「チームの誰かに説明するとしたら、どう言うか？」**
""")
                if st.button("閉じる", use_container_width=True):
                    st.session_state["show_wall_hint"] = False
                    st.rerun()
            wall_hint_dialog()
        prev_krs = master.get("key_results", [{"id": f"kr{i+1}", "label": KR_LABELS[i], "text": "", "walls": []} for i in range(3)])
        while len(prev_krs) < 3:
            i = len(prev_krs)
            prev_krs.append({"id": f"kr{i+1}", "label": KR_LABELS[i], "text": "", "walls": []})

        # 壁リストをsession_stateで管理（初回のみmasterから読み込む）
        for i, kr in enumerate(prev_krs[:3]):
            sk = f"strategy_walls_{i}"
            if sk not in st.session_state:
                walls = kr.get("walls", [])
                st.session_state[sk] = walls if walls else [""]

        kr_texts = []
        kr_walls = []
        for i, kr in enumerate(prev_krs[:3]):
            sk = f"strategy_walls_{i}"
            with st.container(border=True):
                cc, ci = st.columns([0.12, 0.88])
                with cc:
                    st.markdown(f'<div style="background:{KR_COLORS[i]};color:#fff;font-size:.75rem;font-weight:600;border-radius:20px;padding:5px 8px;text-align:center;margin-top:6px;">{KR_LABELS[i]}</div>', unsafe_allow_html=True)
                with ci:
                    val = st.text_input(f"KR{i+1}", value=kr.get("text",""),
                        placeholder="例）月次NPSを 40 以上に引き上げる",
                        key=f"strategy_kr_{i}", label_visibility="collapsed")
                kr_texts.append(val)

                # 壁の入力
                st.markdown(
                    f'<div style="font-size:.75rem;font-weight:500;color:#F39C12;'
                    f'margin-top:.5rem;margin-bottom:.3rem;">🔍 壁（課題）─ チームで合意した内容を入力</div>',
                    unsafe_allow_html=True,
                )
                walls_this_kr = []
                for wi in range(len(st.session_state[sk])):
                    wc, wd = st.columns([0.9, 0.1])
                    with wc:
                        w = st.text_input(
                            f"wall_{i}_{wi}",
                            value=st.session_state[sk][wi],
                            placeholder=f"壁{wi+1}：例）提案資料の訴求力が弱い",
                            label_visibility="collapsed",
                            key=f"strategy_wall_{i}_{wi}",
                        )
                        walls_this_kr.append(w)
                        st.session_state[sk][wi] = w
                    with wd:
                        if wi > 0 and st.button("✕", key=f"del_wall_{i}_{wi}", help="削除"):
                            st.session_state[sk].pop(wi)
                            st.rerun()

                if len(st.session_state[sk]) < MAX_WALLS:
                    if st.button("＋ 壁を追加", key=f"add_wall_{i}"):
                        st.session_state[sk].append("")
                        st.rerun()
                else:
                    st.caption(f"壁は最大 {MAX_WALLS} 件まで")

                kr_walls.append([w for w in walls_this_kr if w.strip()])

    st.markdown("---")
    st.markdown('<div class="g-warn">⚠️ 「確定保存」を押すとOKRがロックされます。全員が納得した上で押してください。</div>', unsafe_allow_html=True)
    col_save, col_hint = st.columns([2, 5])
    with col_save:
        save_btn = st.button("🔒 OKRを確定保存する", type="primary", use_container_width=True)
    with col_hint:
        filled = sum(1 for t in kr_texts if t.strip())
        st.caption(f"Objective：{'✅ 入力済' if objective.strip() else '⬜ 未入力'}　KR入力済：{filled}/3")

    if save_btn:
        if not objective.strip():
            st.error("Objectiveを入力してください。")
        elif not any(t.strip() for t in kr_texts):
            st.error("KRを1つ以上入力してください。")
        elif any(t.strip() and not walls for t, walls in zip(kr_texts, kr_walls)):
            st.error("KRが入力されているものには壁を1つ以上設定してください。")
        else:
            payload = dict(
                quarter    = quarter.strip(),
                objective  = objective.strip(),
                locked     = True,
                set_at     = datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
                key_results= [
                    {"id": f"kr{i+1}", "label": KR_LABELS[i],
                     "text": t.strip(), "walls": kr_walls[i]}
                    for i, t in enumerate(kr_texts) if t.strip()
                ],
            )
            if io_save_master(payload):
                st.session_state.cached_master = payload
                st.toast("🌟 OKRを確定しました！北極星が輝きました。", icon="🌟")
                st.success("✅ 確定保存しました。ページを更新すると北極星バナーに反映されます。")

    # ── データ管理（管理者専用） ───────────────────────────────────────────
    st.markdown("---")
    with st.expander("🗑️ データ管理（テストデータの削除など）"):
        st.markdown('<div class="g-warn">⚠️ 削除したデータは元に戻せません。必ず確認してから実行してください。</div>', unsafe_allow_html=True)

        month_str = st.session_state.month_str
        month_label = datetime.date.fromisoformat(month_str+"-01").strftime("%Y年%m月")

        # 対象月のプラン一覧を表示
        plans = io_list_plans(month_str)
        if not plans:
            st.info(f"📭 {month_label} の保存データはありません。")
        else:
            st.markdown(f"**{month_label} の保存データ一覧**")
            for plan in plans:
                m        = plan.get("member","不明")
                saved_at = plan.get("saved_at","")
                col_info, col_del = st.columns([4, 1])
                with col_info:
                    st.markdown(
                        f'<div style="background:var(--color-background-secondary);'
                        f'border-radius:8px;padding:.5rem .85rem;margin-bottom:.3rem;">'
                        f'<span style="font-weight:600;">{m}</span>'
                        f'<span style="font-size:.75rem;color:var(--color-text-secondary);margin-left:10px;">保存日時：{saved_at}</span>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                with col_del:
                    if st.button("削除", key=f"del_plan_{m}", type="secondary"):
                        io_delete_plan(month_str, m)
                        st.toast(f"🗑️ {m} のデータを削除しました。", icon="🗑️")
                        st.session_state.team_data = None
                        st.rerun()

        st.markdown("---")
        st.markdown("**全データをまとめて削除する**")
        if st.button("🗑️ この月の全データを削除する", type="secondary", use_container_width=True):
            for plan in plans:
                m_ = plan.get("member","")
                if m_:
                    io_delete_plan(month_str, m_)
            st.toast(f"🗑️ {month_label} の全データを削除しました。", icon="🗑️")
            st.session_state.team_data = None
            st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# PLAN タブ
# ══════════════════════════════════════════════════════════════════════════════

def render_plan(master: dict):
    month_str = st.session_state.month_str
    krs       = master.get("key_results", [])
    step      = st.session_state.plan_step
    kr_idx    = st.session_state.plan_kr_idx

    if not krs or not any(kr.get("text") for kr in krs):
        st.warning("OKRがまだ設定されていません。先に「STRATEGY」タブで四半期OKRを確定してください。")
        return

    if not any(kr.get("walls") for kr in krs):
        st.warning("壁がまだ設定されていません。先に「STRATEGY」タブでKRごとの壁を設定してください。")
        return

    draft_key = f"draft_{month_str}_team"
    if draft_key not in st.session_state:
        saved = io_get_plan(month_str)
        if saved:
            loaded = {}
            for item in saved.get("items", []):
                if "wall_actions" in item:
                    loaded[item["kr_id"]] = {"wall_actions": item["wall_actions"]}
                elif "issues" in item:
                    kr_walls = next((kr.get("walls",[]) for kr in krs if kr["id"]==item["kr_id"]), [])
                    wall_actions = []
                    for wi, wall_text in enumerate(kr_walls):
                        old_iss = item["issues"][wi] if wi < len(item["issues"]) else {}
                        wall_actions.append({
                            "wall_text": wall_text,
                            "actions":   old_iss.get("actions", [{"text":"","start":"","end":""}]),
                        })
                    loaded[item["kr_id"]] = {"wall_actions": wall_actions}
            st.session_state[draft_key] = loaded
            st.success(f"✅ 前回保存（{saved.get('saved_at','')}）を読み込みました。")
        else:
            st.session_state[draft_key] = {}
            st.info("📋 今月のチームプランはまだありません。KRを選んでアクションを入力しましょう。")

    draft: dict = st.session_state[draft_key]

    # 優先度データ（セッションキャッシュ。DASHBOARDで確定後に更新される）
    if "priorities" not in st.session_state:
        st.session_state.priorities = io_get_priorities(month_str)
    priorities = st.session_state.priorities

    render_progress(step)

    # ── STEP 0: KR選択 ────────────────────────────────────────────────────
    if step == 0:
        st.markdown('<div class="g-info">チームで取り組むKRを選んでください。全員で画面を見ながら議論してアクションを決めましょう。</div>', unsafe_allow_html=True)
        st.markdown("")
        for i, kr in enumerate(krs):
            if not kr.get("walls"):
                continue
            kd          = draft.get(kr["id"], {})
            wall_actions = kd.get("wall_actions", [])
            done = any(
                any(a.get("text","").strip() for a in wa.get("actions",[]))
                for wa in wall_actions
            )
            col_text, col_btn = st.columns([6, 1])
            with col_text:
                walls_list = "".join(
                    f'<div style="display:flex;align-items:flex-start;gap:5px;margin-top:.2rem;">'
                    f'<span style="font-size:.7rem;font-weight:600;color:#F39C12;flex-shrink:0;">壁{wi+1}</span>'
                    f'<span style="font-size:.72rem;color:var(--color-text-secondary);line-height:1.4;">{w}</span>'
                    f'</div>'
                    for wi, w in enumerate(kr.get("walls",[]))
                )
                st.markdown(
                    f'<div style="background:{KR_COLORS[i]}18;border:1.5px solid {KR_COLORS[i]};'
                    f'border-radius:10px;padding:.65rem 1rem;margin-bottom:.4rem;">'
                    f'<div style="display:flex;align-items:center;justify-content:space-between;">'
                    f'<div><span style="background:{KR_COLORS[i]};color:#fff;font-size:.7rem;font-weight:700;padding:2px 8px;border-radius:20px;">{kr["label"]}</span>'
                    f'<span style="font-size:.88rem;font-weight:500;color:{KR_COLORS[i]};margin-left:8px;">{kr["text"]}</span></div>'
                    f'<span style="font-size:.72rem;color:{"#196F3D" if done else "#7F8C8D"};">{"✅ 入力済" if done else "⬜ 未入力"}</span>'
                    f'</div>{walls_list}</div>', unsafe_allow_html=True,
                )
            with col_btn:
                if st.button("選ぶ", key=f"sel_kr_{i}", use_container_width=True):
                    st.session_state.plan_kr_idx = i
                    walls = kr.get("walls", [])
                    saved_wa = kd.get("wall_actions", [])
                    merged = []
                    for wi, wall_text in enumerate(walls):
                        existing = next((wa for wa in saved_wa if wa.get("wall_text")==wall_text), None)
                        if existing:
                            merged.append(existing)
                        else:
                            merged.append({"wall_text": wall_text, "actions": [{"text":"","start":"","end":""}]})
                    st.session_state.plan_actions = merged
                    st.session_state.plan_step    = 1
                    st.rerun()

    # ── STEP 1: アクション入力 ────────────────────────────────────────────
    elif step == 1:
        kr          = krs[kr_idx]
        wall_actions = st.session_state.plan_actions

        col_form, col_tree = st.columns([3, 2], gap="medium")

        with col_form:
            st.markdown(
                f'<div style="background:{KR_COLORS[kr_idx]}18;border-left:4px solid {KR_COLORS[kr_idx]};'
                f'padding:.5rem .85rem;border-radius:0 8px 8px 0;font-weight:600;font-size:.85rem;'
                f'margin-bottom:.85rem;color:{KR_COLORS[kr_idx]};">{kr["label"]}：{kr["text"]}</div>',
                unsafe_allow_html=True,
            )
            st.markdown('<div class="g-info">チームで合意した<b>壁（課題）</b>ごとに、今月のアクションを入力してください。担当者は優先度設定画面で割り当てます。</div>', unsafe_allow_html=True)

            if st.button("💡 アクションの考え方ヒント", use_container_width=False):
                st.session_state["show_action_hint_plan"] = True
            if st.session_state.get("show_action_hint_plan"):
                @st.dialog("💡 アクションの考え方ヒント")
                def action_hint_plan_dialog():
                    st.markdown("### アクションとは？")
                    st.markdown("この壁を乗り越えるために、**チームが今月やること**を自由に書いてください。まずはアイデアをどんどん出しましょう。後で優先度と担当者を決めます。")
                    st.markdown("---")
                    st.markdown("### ✅ 書き方の例")
                    st.markdown("""
- 「顧客インタビューを実施する」
- 「提案資料に導入事例を追加する」
- 「6/15までにインタビュー10件実施してSlackで共有する」
""")
                    st.markdown("### 💡 ポイント")
                    st.markdown("最初は気軽に書いてOKです。日付や件数は後から追加できます。")
                    if st.button("閉じる", use_container_width=True):
                        st.session_state["show_action_hint_plan"] = False
                        st.rerun()
                action_hint_plan_dialog()

            for ii, wa in enumerate(wall_actions):
                with st.container(border=True):
                    st.markdown(
                        f'<div style="display:flex;align-items:flex-start;gap:7px;margin-bottom:.6rem;'
                        f'background:#FEF9E7;border-radius:8px;padding:.55rem .75rem;">'
                        f'<div style="width:22px;height:22px;border-radius:50%;background:#F39C12;'
                        f'display:flex;align-items:center;justify-content:center;font-size:.7rem;'
                        f'font-weight:600;color:#fff;flex-shrink:0;margin-top:1px;">壁{ii+1}</div>'
                        f'<span style="font-size:.85rem;font-weight:600;color:#7D6608;line-height:1.5;">{wa["wall_text"]}</span>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                    st.markdown(
                        '<div style="font-size:10px;color:var(--color-text-secondary);'
                        'margin-bottom:.3rem;font-weight:500;">⚡ この壁に対するアクション（自由記述）</div>',
                        unsafe_allow_html=True,
                    )

                    actions = wa.get("actions", [{"text":"","start":"","end":""}])
                    wall_actions[ii]["actions"] = actions
                    actions_to_delete = []

                    for ia, action in enumerate(actions):
                        with st.container(border=False):
                            col_ahdr, col_adel = st.columns([8, 1])
                            with col_ahdr:
                                ak  = f"team__{kr['id']}__{ii}__{ia}"
                                saved_pri = priorities.get(ak, {})
                                pri = saved_pri.get("priority","") if isinstance(saved_pri, dict) else saved_pri
                                pri_colors = {"高": "#E74C3C", "中": "#F39C12", "低": "#27AE60"}
                                pri_badge  = (
                                    f'<span style="background:{pri_colors[pri]};color:#fff;'
                                    f'font-size:.65rem;font-weight:600;padding:1px 7px;'
                                    f'border-radius:20px;margin-left:6px;">{pri}</span>'
                                ) if pri in pri_colors else ""
                                # 担当者バッジ
                                assignee = saved_pri.get("assignee","") if isinstance(saved_pri, dict) else ""
                                assignee_badge = (
                                    f'<span style="font-size:.65rem;color:var(--color-text-secondary);margin-left:6px;">👤{assignee}</span>'
                                ) if assignee else ""
                                st.markdown(
                                    f'<div style="display:flex;align-items:center;gap:6px;'
                                    f'background:var(--color-background-secondary);border-radius:6px;padding:.3rem .5rem;">'
                                    f'<div style="width:18px;height:18px;border-radius:50%;background:#1B4F72;'
                                    f'display:flex;align-items:center;justify-content:center;font-size:.62rem;font-weight:600;color:#fff;">A{ia+1}</div>'
                                    f'<span style="font-size:.8rem;font-weight:500;color:var(--color-text-primary);">アクション {ia+1}</span>'
                                    f'{pri_badge}{assignee_badge}</div>',
                                    unsafe_allow_html=True,
                                )
                            with col_adel:
                                if st.button("✕", key=f"del_action_{ii}_{ia}", help="削除"):
                                    actions_to_delete.append(ia)

                            actions[ia]["text"] = st.text_area(
                                f"act_{ii}_{ia}",
                                value=action.get("text",""), height=65,
                                placeholder="例）顧客インタビューを実施する、提案資料に導入事例を追加する…",
                                label_visibility="collapsed", key=f"act_txt_{ii}_{ia}",
                            )
                            use_date = st.checkbox(
                                "日付を設定する（任意）",
                                value=bool(action.get("start","")),
                                key=f"use_date_{ii}_{ia}",
                            )
                            if use_date:
                                c1, c2 = st.columns(2)
                                with c1:
                                    try:
                                        sv = datetime.date.fromisoformat(action.get("start","") or datetime.date.today().isoformat())
                                    except ValueError:
                                        sv = datetime.date.today()
                                    actions[ia]["start"] = st.date_input("開始日", value=sv, key=f"start_{ii}_{ia}").isoformat()
                                with c2:
                                    try:
                                        ev = datetime.date.fromisoformat(action.get("end","") or (datetime.date.today()+datetime.timedelta(days=30)).isoformat())
                                    except ValueError:
                                        ev = datetime.date.today()+datetime.timedelta(days=30)
                                    actions[ia]["end"] = st.date_input("終了日", value=ev, key=f"end_{ii}_{ia}").isoformat()
                            else:
                                actions[ia]["start"] = ""
                                actions[ia]["end"]   = ""

                    for ia in sorted(actions_to_delete, reverse=True):
                        if len(actions) > 1:
                            actions.pop(ia)
                        else:
                            st.warning("アクションは最低1件必要です。")
                    if actions_to_delete:
                        st.rerun()

                    if len(actions) < MAX_ACTIONS:
                        if st.button("＋ アクションを追加", key=f"add_action_{ii}", use_container_width=True):
                            actions.append({"text":"","start":"","end":""}); st.rerun()
                    else:
                        st.caption(f"アクションは最大 {MAX_ACTIONS} 件まで")

            st.markdown("---")
            c_back, c_other, c_save = st.columns([1, 1.5, 2])
            with c_back:
                if st.button("← KR選択に戻る", use_container_width=True):
                    st.session_state.plan_step = 0; st.rerun()
            with c_other:
                if st.button("別のKRも入力する", use_container_width=True):
                    draft[krs[kr_idx]["id"]] = {"wall_actions": wall_actions}
                    st.session_state.plan_step = 0; st.rerun()
            with c_save:
                has_valid = any(
                    any(a.get("text","").strip() for a in wa.get("actions",[]))
                    for wa in wall_actions
                )
                if st.button("💾 保存する", type="primary", use_container_width=True, disabled=not has_valid):
                    draft[krs[kr_idx]["id"]] = {"wall_actions": wall_actions}
                    payload = dict(
                        month    = month_str,
                        saved_at = datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
                        items    = [
                            dict(
                                kr_id       = kr_["id"],
                                kr_label    = kr_["label"],
                                kr_text     = kr_["text"],
                                wall_actions= [
                                    {
                                        "wall_text": wa.get("wall_text",""),
                                        "actions":   [
                                            {"text": a.get("text",""), "start": a.get("start",""), "end": a.get("end","")}
                                            for a in wa.get("actions",[]) if a.get("text","").strip()
                                        ],
                                    }
                                    for wa in draft.get(kr_["id"],{}).get("wall_actions",[])
                                ],
                            )
                            for kr_ in krs
                            if draft.get(kr_["id"],{}).get("wall_actions")
                        ],
                    )
                    if io_save_plan(month_str, payload):
                        st.toast("✅ チームプランを保存しました！", icon="🎉")
                        st.success("保存完了！「DASHBOARD」タブで確認できます。")
                        del st.session_state[draft_key]
                        st.session_state.plan_step = 0; st.rerun()

        with col_tree:
            render_logic_tree(master, kr_idx, wall_actions, {"main": "#1B4F72"})

def build_gantt(all_plans: list[dict]) -> Any | None:
    rows = []
    priorities = {}
    if all_plans:
        # 担当者情報を優先度データから取得
        try:
            month_str = all_plans[0].get("month","")
            if month_str:
                priorities = io_get_priorities(month_str)
        except Exception:
            pass

    for plan in all_plans:
        for item in plan.get("items",[]):
            for ii, wa in enumerate(item.get("wall_actions",[])):
                for ia, action in enumerate(wa.get("actions",[])):
                    if not action.get("text") or not action.get("start") or not action.get("end"):
                        continue
                    ak       = f"team__{item['kr_id']}__{ii}__{ia}"
                    saved_p  = priorities.get(ak, {})
                    assignee = saved_p.get("assignee","チーム") if isinstance(saved_p, dict) else "チーム"
                    rows.append(dict(
                        member = assignee,
                        label  = f'{assignee}｜{item["kr_label"]}',
                        action = action["text"],
                        start  = pd.Timestamp(action["start"]),
                        end    = pd.Timestamp(action["end"]) + pd.Timedelta(days=1),
                    ))
    if not rows:
        return None
    df = pd.DataFrame(rows)
    fig = px.timeline(df, x_start="start", x_end="end", y="label", color="member",
        color_discrete_map={m: mpal(m)["main"] for m in MEMBERS},
        hover_data={"action": True, "start": False, "end": False},
        labels={"label": ""})
    fig.update_yaxes(autorange="reversed", tickfont=dict(size=11))
    fig.update_xaxes(title="", tickfont=dict(size=11))
    fig.update_layout(
        plot_bgcolor="white", paper_bgcolor="white",
        font=dict(family="Calibri, Arial", size=12, color="#1B4F72"),
        legend=dict(title="担当者", orientation="h", yanchor="bottom", y=1.02, xanchor="left", x=0),
        margin=dict(l=10, r=10, t=44, b=10),
        height=max(300, len(rows)*44+80))
    fig.update_traces(marker_line_width=0, opacity=0.88)
    return fig


# ══════════════════════════════════════════════════════════════════════════════
# DASHBOARD タブ
# ══════════════════════════════════════════════════════════════════════════════

def render_dashboard(master: dict):
    month_str   = st.session_state.month_str
    month_label = datetime.date.fromisoformat(month_str+"-01").strftime("%Y年%m月")

    fc, _ = st.columns([2, 6])
    with fc:
        if st.button("🔄 データを読み込む", use_container_width=True, type="primary"):
            st.session_state.team_data = None

    if st.session_state.team_data is None:
        st.session_state.team_data = io_list_plans(month_str)

    all_plans: list[dict] = st.session_state.team_data or []

    if not all_plans:
        st.info("まだチームプランが保存されていません。PLANタブで入力・保存してください。")
        return

    plan = all_plans[0]  # チームプランは1件

    # 統計
    total_actions = sum(
        len(wa.get("actions",[]))
        for item in plan.get("items",[])
        for wa in item.get("wall_actions",[])
    )
    all_starts = [
        a["start"] for item in plan.get("items",[])
        for wa in item.get("wall_actions",[]) for a in wa.get("actions",[]) if a.get("start")
    ]
    all_ends = [
        a["end"] for item in plan.get("items",[])
        for wa in item.get("wall_actions",[]) for a in wa.get("actions",[]) if a.get("end")
    ]

    st.markdown('<div class="stat-row">', unsafe_allow_html=True)
    for n, lbl in [
        (total_actions, "合計アクション数"),
        (min(all_starts) if all_starts else "-", "最早 開始日"),
        (max(all_ends)   if all_ends   else "-", "最遅 終了日"),
        (plan.get("saved_at",""), "最終更新"),
    ]:
        st.markdown(f'<div class="stat"><div class="stat-n" style="font-size:1.1rem;">{n}</div><div class="stat-l">{lbl}</div></div>', unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

    # フィードバックパネル
    with st.expander("🔍 レビューの着眼点", expanded=False):
        st.markdown("""
<div class="fb-panel">
<h4>📌 チェックリスト</h4>
<div class="fb-item"><span class="fb-ico">🔥</span><div><b>野心度</b>　70〜80%達成が理想。簡単すぎると成長が止まる。</div></div>
<div class="fb-item"><span class="fb-ico">🔗</span><div><b>ロジック</b>　「壁→アクション」の因果関係は通っているか？</div></div>
<div class="fb-item"><span class="fb-ico">⚡</span><div><b>リソース</b>　全アクション合計で1ヶ月以内に現実的にやり切れるか？</div></div>
<div class="fb-item"><span class="fb-ico">🎯</span><div><b>整合性</b>　全アクションが同じObjectiveを向いているか？</div></div>
</div>
""", unsafe_allow_html=True)

    # チームプランカード
    st.markdown("### 📋 チーム アクション一覧")
    priorities = io_get_priorities(month_str)
    for item in plan.get("items",[]):
        kr_idx_ = next((i for i,kr in enumerate(master.get("key_results",[])) if kr.get("id")==item.get("kr_id")), 0)
        kr_col  = KR_COLORS[kr_idx_ % len(KR_COLORS)]
        walls_html = ""
        for ii, wa in enumerate(item.get("wall_actions",[])):
            a_html = ""
            for ai, a in enumerate(wa.get("actions",[])):
                if not a.get("text","").strip():
                    continue
                ak       = f"team__{item['kr_id']}__{ii}__{ai}"
                saved_p  = priorities.get(ak, {})
                pri      = saved_p.get("priority","") if isinstance(saved_p, dict) else saved_p
                assignee = saved_p.get("assignee","") if isinstance(saved_p, dict) else ""
                pri_colors = {"高": "#E74C3C", "中": "#F39C12", "低": "#27AE60"}
                pri_badge = f'<span style="background:{pri_colors[pri]};color:#fff;font-size:.65rem;font-weight:600;padding:1px 6px;border-radius:20px;margin-left:4px;">{pri}</span>' if pri in pri_colors else ""
                assignee_badge = f'<span style="font-size:.7rem;color:var(--color-text-secondary);margin-left:6px;">👤 {assignee}</span>' if assignee else ""
                date_text = f'{a.get("start","")} → {a.get("end","")}' if a.get("start") or a.get("end") else ""
                a_html += (
                    f'<div class="action-list-item">'
                    f'<div class="a-num" style="background:#1B4F72;">A{ai+1}</div>'
                    f'<div style="flex:1;"><div style="display:flex;align-items:center;flex-wrap:wrap;">'
                    f'<span style="font-size:.8rem;color:var(--color-text-primary);">{a.get("text","")}</span>'
                    f'{pri_badge}{assignee_badge}</div>'
                    f'{"<div style=\'font-size:.7rem;color:var(--color-text-secondary);\'>" + date_text + "</div>" if date_text else ""}'
                    f'</div></div>'
                )
            walls_html += (
                f'<div style="background:var(--color-background-primary);border-radius:6px;'
                f'padding:.35rem .6rem;margin-bottom:.35rem;">'
                f'<div style="display:flex;align-items:center;gap:5px;margin-bottom:.25rem;">'
                f'<div style="width:18px;height:18px;border-radius:50%;background:#F39C12;'
                f'display:flex;align-items:center;justify-content:center;font-size:.62rem;font-weight:600;color:#fff;">壁{ii+1}</div>'
                f'<span style="font-size:.78rem;font-weight:500;color:#7D6608;">{wa["wall_text"]}</span>'
                f'</div>{a_html}</div>'
            )
        st.markdown(
            f'<div class="kr-block">'
            f'<div style="display:flex;align-items:center;gap:6px;margin-bottom:.35rem;">'
            f'<span style="background:{kr_col};color:#fff;font-size:.68rem;font-weight:700;padding:1px 8px;border-radius:20px;">{item["kr_label"]}</span>'
            f'<span style="font-size:.78rem;font-weight:500;color:var(--color-text-primary);">{item.get("kr_text","")[:45]}</span>'
            f'</div>{walls_html}</div>', unsafe_allow_html=True,
        )

    # ── 優先度設定（管理者専用） ──────────────────────────────────────────
    st.markdown("### 🎯 アクション優先度設定")
    if not st.session_state.get("admin_auth"):
        st.markdown('<div class="g-info">管理者ログイン後に優先度を設定できます。STRATEGYタブで認証してください。</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="g-info">効果（1〜5）と手間（1〜5）を入力すると優先度が自動提案されます。内容を確認して修正後、「✅ 優先度を確定する」を押してください。</div>', unsafe_allow_html=True)

        # スコア→優先度変換
        SIZE_MAP = {"小": 1, "中": 2, "大": 3, "": 0}
        SIZE_OPTIONS = ["", "小", "中", "大"]

        def calc_priority(effect: str, effort: str) -> str:
            e1 = SIZE_MAP.get(effect, 0)
            e2 = SIZE_MAP.get(effort, 0)
            if e1 == 0 or e2 == 0:
                return "未設定"
            score = e1 / e2
            if score >= 1.5:
                return "高"
            elif score >= 0.7:
                return "中"
            else:
                return "低"

        PRIORITY_OPTIONS = ["未設定", "高", "中", "低"]
        PRIORITY_COLORS  = {"高": "#E74C3C", "中": "#F39C12", "低": "#27AE60", "未設定": "#95A5A6"}

        # 優先度データ読み込み
        if "priorities" not in st.session_state:
            st.session_state.priorities = io_get_priorities(month_str)
        priorities: dict = st.session_state.priorities

        if not all_plans:
            st.warning("先に「🔄 データを読み込む」ボタンを押してください。")
        else:
            sort_by_priority = st.checkbox("自動提案の優先度でソートする", value=False, key="sort_priority")
    
            PRIORITY_ORDER = {"高": 0, "中": 1, "低": 2, "未設定": 3}
    
            # 全アクションをフラット化してソート
            all_action_rows = []
            for plan in all_plans:
                for item in plan.get("items",[]):
                    for ii, wa in enumerate(item.get("wall_actions",[])):
                        for ia, a in enumerate(wa.get("actions",[])):
                            if not a.get("text","").strip():
                                continue
                            ak    = f"team__{item['kr_id']}__{ii}__{ia}"
                            saved = priorities.get(ak, {})
                            if isinstance(saved, str):
                                saved = {"effect": "", "effort": "", "priority": saved}
                            auto  = calc_priority(saved.get("effect",""), saved.get("effort",""))
                            all_action_rows.append({
                                "item": item, "ii": ii, "ia": ia,
                                "wa": wa, "a": a, "ak": ak,
                                "saved": saved, "auto": auto,
                            })
    
            if sort_by_priority:
                all_action_rows.sort(key=lambda r: PRIORITY_ORDER.get(r["auto"], 3))
    
            # ヘッダー
            h_cols = st.columns([2.5, 2, 1.1, 1.1, 1.4, 1.6, 1.6])
            for col, label in zip(h_cols, ["アクション", "壁", "効果", "手間", "自動提案", "優先度（確定）", "担当者"]):
                with col:
                    st.markdown(
                        f'<div style="font-size:.72rem;font-weight:600;color:var(--color-text-secondary);'
                        f'padding:.2rem 0;border-bottom:1.5px solid var(--color-border-secondary);">{label}</div>',
                        unsafe_allow_html=True,
                    )
    
            cur_member_label = None
            for row in all_action_rows:
                item = row["item"]
                ii, ia, wa, a = row["ii"], row["ia"], row["wa"], row["a"]
                ak, saved, auto = row["ak"], row["saved"], row["auto"]
    
                color = PRIORITY_COLORS.get(auto, "#95A5A6")
    
                c_act, c_wall, c_eff, c_eft, c_auto, c_pri, c_mem = st.columns([2.5, 2, 1.1, 1.1, 1.4, 1.6, 1.6])
    
                with c_act:
                    st.markdown(
                        f'<div style="font-size:.78rem;color:var(--color-text-primary);'
                        f'padding:.3rem 0;line-height:1.4;">{a["text"]}</div>',
                        unsafe_allow_html=True,
                    )
                with c_wall:
                    st.markdown(
                        f'<div style="font-size:.7rem;color:#7D6608;padding:.3rem 0;">'
                        f'壁{ii+1}：{wa["wall_text"][:18]}{"…" if len(wa["wall_text"])>18 else ""}</div>',
                        unsafe_allow_html=True,
                    )
                with c_eff:
                    eff_val = saved.get("effect", "")
                    if eff_val not in SIZE_OPTIONS:
                        eff_val = ""
                    effect = st.selectbox(
                        "効果", SIZE_OPTIONS,
                        index=SIZE_OPTIONS.index(eff_val),
                        key=f"effect_{ak}",
                        label_visibility="collapsed",
                    )
                with c_eft:
                    eft_val = saved.get("effort", "")
                    if eft_val not in SIZE_OPTIONS:
                        eft_val = ""
                    effort = st.selectbox(
                        "手間", SIZE_OPTIONS,
                        index=SIZE_OPTIONS.index(eft_val),
                        key=f"effort_{ak}",
                        label_visibility="collapsed",
                    )
                with c_auto:
                    auto_now = calc_priority(effect, effort)
                    color    = PRIORITY_COLORS.get(auto_now, "#95A5A6")
                    st.markdown(
                        f'<div style="background:{color};color:#fff;font-size:.72rem;'
                        f'font-weight:600;padding:3px 10px;border-radius:20px;'
                        f'text-align:center;margin-top:4px;">{auto_now}</div>',
                        unsafe_allow_html=True,
                    )
                with c_pri:
                    current_pri = saved.get("priority", auto_now)
                    if current_pri not in PRIORITY_OPTIONS:
                        current_pri = auto_now if auto_now in PRIORITY_OPTIONS else "未設定"
                    final_pri = st.selectbox(
                        "優先度",
                        PRIORITY_OPTIONS,
                        index=PRIORITY_OPTIONS.index(current_pri),
                        key=f"pri_{ak}",
                        label_visibility="collapsed",
                    )
                with c_mem:
                    current_assignee = saved.get("assignee", MEMBERS[0])
                    if current_assignee not in MEMBERS:
                        current_assignee = MEMBERS[0]
                    assignee = st.selectbox(
                        "担当者",
                        MEMBERS,
                        index=MEMBERS.index(current_assignee),
                        key=f"assignee_{ak}",
                        label_visibility="collapsed",
                    )
    
                priorities[ak] = {
                    "effect":   effect,
                    "effort":   effort,
                    "auto":     auto_now,
                    "priority": final_pri,
                    "assignee": assignee,
                }
    
            st.markdown("")
            if st.button("✅ 優先度を確定する", type="primary", use_container_width=False):
                if io_save_priorities(month_str, priorities):
                    st.session_state.priorities = priorities
                    st.toast("✅ 優先度を確定しました！", icon="🎯")
                    st.rerun()

    # ガントチャート
    st.markdown("### 📊 統合ガントチャート")
    fig = build_gantt(all_plans)
    if fig:
        st.plotly_chart(fig, use_container_width=True)
    else:
        st.info("アクションと日付が入力されたデータがあるとガントチャートが表示されます。")

    # PPTX
    st.markdown("---")
    st.markdown("### 📥 PPTX ダウンロード")
    gen_c, hint_c = st.columns([2, 5])
    with gen_c:
        gen_btn = st.button("🚀 PPTXを生成する", type="primary", use_container_width=True)
    with hint_c:
        st.caption(f"構成：表紙 + チームOKR + メンバー別詳細({len(all_plans)}枚) + ガントチャート ＝ 計{2+len(all_plans)+1}枚")

    if gen_btn:
        with st.spinner("PPTX生成中…"):
            try:
                gantt_png  = fig.to_image(format="png", width=1400, scale=2) if fig else None
                pptx_bytes = build_pptx(CFG["team_name"], month_label, master, all_plans, gantt_png)
                fname = f"OKR_{CFG['team_name']}_{month_str}.pptx"
                st.toast("🎉 PPTXの生成が完了しました！", icon="📊")
                st.download_button(
                    label="⬇️ PPTXをダウンロード", data=pptx_bytes, file_name=fname,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )
            except Exception:
                st.error("PPTX生成中にエラーが発生しました。")
                st.code(traceback.format_exc())


# ══════════════════════════════════════════════════════════════════════════════
# PPTX 生成
# ══════════════════════════════════════════════════════════════════════════════

def _rgb(h):
    h=h.lstrip("#"); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def _rect(sl,x,y,w,h,fill,border=None,bw=0.5):
    s=sl.shapes.add_shape(1,x,y,w,h); s.fill.solid(); s.fill.fore_color.rgb=fill
    if border: s.line.color.rgb=border; s.line.width=Pt(bw)
    else: s.line.fill.background()
    return s

def _txt(sl,text,x,y,w,h,sz,bold=False,color=None,align=PP_ALIGN.LEFT,italic=False):
    color=color or RGBColor(0xFF,0xFF,0xFF)
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.size=Pt(sz); r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=color; r.font.name="Calibri"
    return tb

# ══════════════════════════════════════════════════════════════════════════════
# タスク起票タブ
# ══════════════════════════════════════════════════════════════════════════════

def render_task_ticket(master: dict, month_str: str):
    st.markdown('<div class="g-info">優先度「高」「中」のアクションを選んで、一問一答形式でタスクを具体化します。最後にBacklog用のマークダウンが生成されます。</div>', unsafe_allow_html=True)

    priorities = io_get_priorities(month_str)
    all_plans  = io_list_plans(month_str)

    if not all_plans:
        st.warning("まだ誰も計画を提出していません。PLANタブで入力・保存してください。")
        return
    if not priorities:
        st.warning("まだ優先度が設定されていません。DASHBOARDタブで優先度を確定してください。")
        return

    TARGET_PRI = {"高", "中"}
    action_options = []
    for plan in all_plans:
        for item in plan.get("items",[]):
            for ii, wa in enumerate(item.get("wall_actions",[])):
                for ia, a in enumerate(wa.get("actions",[])):
                    if not a.get("text","").strip():
                        continue
                    ak    = f"team__{item['kr_id']}__{ii}__{ia}"
                    saved = priorities.get(ak, {})
                    pri      = saved.get("priority","未設定") if isinstance(saved, dict) else saved
                    assignee = saved.get("assignee", MEMBERS[0]) if isinstance(saved, dict) else MEMBERS[0]
                    if pri in TARGET_PRI:
                        action_options.append({
                            "label":    f"【{pri}】{assignee}｜{item['kr_label']}｜{a['text'][:30]}{'…' if len(a['text'])>30 else ''}",
                            "assignee": assignee,
                            "kr_label": item["kr_label"],
                            "kr_text":  item["kr_text"],
                            "wall":     wa["wall_text"],
                            "action":   a["text"],
                            "priority": pri,
                            "ak":       ak,
                        })

    if not action_options:
        st.warning("優先度「高」「中」のアクションがありません。DASHBOARDタブで優先度を設定してください。")
        return

    st.markdown("### STEP 1：アクションを選ぶ")
    labels   = [o["label"] for o in action_options]
    selected = st.selectbox("対象アクション", labels, label_visibility="collapsed")
    sel      = action_options[labels.index(selected)]

    obj = master.get("objective","")
    st.markdown(f"""
<div class="ltree" style="margin:.8rem 0;">
  <div class="ltree-hdr">🌲 ロジックツリー</div>
  <div class="lt-row"><div class="lt-ico" style="background:#1B4F72;">O</div>
    <div class="lt-body" style="font-weight:600;">{obj}</div></div>
  <div class="lt-indent">
    <div class="lt-row"><div class="lt-ico" style="background:#2E86C1;">KR</div>
      <div class="lt-body" style="font-weight:600;">{sel["kr_label"]}：{sel["kr_text"]}</div></div>
    <div class="lt-indent">
      <div class="lt-row"><div class="lt-ico" style="background:#F39C12;">壁</div>
        <div class="lt-body" style="font-weight:600;">{sel["wall"]}</div></div>
      <div class="lt-indent">
        <div class="lt-row"><div class="lt-ico" style="background:#1D6A45;">A</div>
          <div class="lt-body" style="font-weight:600;">{sel["action"]}</div></div>
      </div>
    </div>
  </div>
</div>
""", unsafe_allow_html=True)

    st.markdown("### STEP 2：一問一答で詳細を入力")
    st.markdown('<div class="g-info">必須項目（Q1〜Q8）に回答してください。任意項目はスキップできます。</div>', unsafe_allow_html=True)

    ak = sel["ak"]
    sk = f"task_qa_{ak}"
    if sk not in st.session_state:
        st.session_state[sk] = {}
    qa: dict = st.session_state[sk]

    QUESTIONS = [
        ("target",       "Q1. 相手先は誰ですか？",                                True,  "text",   "例）入居者家族、事業者、〇〇部署"),
        ("coordinator",  "Q2. 誰と調整が必要ですか？",                             True,  "text",   "例）営業担当、現場スタッフ、〇〇さん"),
        ("resources",    "Q3. 必要なものは何ですか？（資料・ツール・データなど）",   True,  "text",   "例）インタビューシート、録音ツール、顧客リスト"),
        ("first_action", "Q4. 最初にやることは何ですか？",                          True,  "text",   "例）対象者リストを作成する"),
        ("deadline",     "Q5. 期限はいつですか？",                                  True,  "date",   ""),
        ("assignee",     "Q6. 担当者は誰ですか？",                                  True,  "member", ""),
        ("done_def",     "Q7. 完了したとみなせる状態は何ですか？",                   True,  "text",   "例）インタビュー10件実施しSlackで共有完了"),
        ("next_action",  "Q8. 完了後のネクストアクションは何ですか？",               True,  "text",   "例）結果をまとめてチームに報告、次の施策を検討"),
        ("approver",     "Q9. 承認・確認を取る必要がある人は？（任意）",             False, "text",   "例）上長、〇〇マネジャー"),
        ("share",        "Q10. 完了後に誰にどう共有しますか？（任意）",              False, "text",   "例）Slackで全員に共有、月例MTGで報告"),
        ("risk",         "Q11. うまくいかない可能性があるとしたら何ですか？（任意）", False, "text",   "例）相手の都合が合わない、データが揃わない"),
    ]

    all_required_filled = True
    for key, label, required, input_type, placeholder in QUESTIONS:
        badge = "" if required else ' <span style="color:var(--color-text-secondary);font-size:.75rem;">(任意)</span>'
        st.markdown(
            f'<div style="font-size:.85rem;font-weight:{"600" if required else "400"};'
            f'color:var(--color-text-primary);margin-top:.8rem;">{label}{badge}</div>',
            unsafe_allow_html=True,
        )
        if input_type == "text":
            val = st.text_area(
                label, value=qa.get(key,""), height=75,
                placeholder=placeholder, label_visibility="collapsed",
                key=f"qa_{ak}_{key}",
            )
            qa[key] = val
            if required and not val.strip():
                all_required_filled = False

        elif input_type == "date":
            try:
                dv = datetime.date.fromisoformat(qa.get(key,"") or datetime.date.today().isoformat())
            except ValueError:
                dv = datetime.date.today()
            val = st.date_input(label, value=dv, label_visibility="collapsed", key=f"qa_{ak}_{key}")
            qa[key] = val.isoformat()

        elif input_type == "member":
            default_a = sel["assignee"] if sel["assignee"] in MEMBERS else MEMBERS[0]
            cur = qa.get(key, default_a)
            if cur not in MEMBERS:
                cur = default_a
            val = st.selectbox(label, MEMBERS, index=MEMBERS.index(cur),
                               label_visibility="collapsed", key=f"qa_{ak}_{key}")
            qa[key] = val

    st.markdown("---")
    st.markdown("### STEP 3：Backlog用マークダウンを生成する")

    if not all_required_filled:
        st.warning("⚠️ 必須項目（Q1〜Q8）をすべて入力してください。")
    else:
        optional_lines = ""
        if qa.get("approver","").strip():
            optional_lines += f"- **承認者**：{qa['approver']}\n"
        if qa.get("share","").strip():
            optional_lines += f"- **完了後の共有**：{qa['share']}\n"
        if qa.get("risk","").strip():
            optional_lines += f"- **リスク**：{qa['risk']}\n"

        md = f"""## 【アクション】{sel["action"]}

### ロジック
- **O**：{obj}
- **{sel["kr_label"]}**：{sel["kr_text"]}
- **壁**：{sel["wall"]}
- **アクション**：{sel["action"]}
- **優先度**：{sel["priority"]}

### タスク詳細
- **相手先**：{qa.get("target","")}
- **調整相手**：{qa.get("coordinator","")}
- **必要なもの**：{qa.get("resources","")}
- **最初にやること**：{qa.get("first_action","")}
- **期限**：{qa.get("deadline","")}
- **担当者**：{qa.get("assignee","")}
{optional_lines}
### 完了の定義
{qa.get("done_def","")}

### 完了後のネクストアクション
{qa.get("next_action","")}
"""
        st.markdown('<div class="g-ok">✅ 必須項目がすべて入力されました。マークダウンをコピーしてBacklogに貼り付けてください。</div>', unsafe_allow_html=True)
        st.code(md, language="markdown")
        st.link_button(
            "🔗 Backlogで新規チケットを開く",
            "https://kyuden-ict.backlog.com/add/MIMAMORIOPS",
            use_container_width=True,
        )


def build_pptx(team_name,month_label,master,all_plans,gantt_png):
    C_NAVY=_rgb("#0D1B2A"); C_DEEP=_rgb("#1B4F72"); C_GRN=_rgb("#1D6A45")
    C_GOLD=_rgb("#F39C12"); C_SNOW=_rgb("#F8FAFC"); C_SLAT=_rgb("#334155")
    C_GRAY=_rgb("#95A5A6"); C_WHT=RGBColor(0xFF,0xFF,0xFF)
    W=Inches(13.33); H=Inches(7.5)
    prs=Presentation(); prs.slide_width=W; prs.slide_height=H
    blank=prs.slide_layouts[6]
    obj=master.get("objective","")
    krs=master.get("key_results",[])

    # 表紙
    s1=prs.slides.add_slide(blank)
    _rect(s1,0,0,W,H*.62,C_NAVY); _rect(s1,0,H*.62,W,H*.38,C_SNOW)
    _rect(s1,Inches(.9),Inches(1.0),Inches(.07),Inches(3.2),C_GOLD)
    _txt(s1,"OKR × 業務計画",Inches(1.1),Inches(.9),Inches(10),Inches(1.1),sz=44,bold=True,color=C_WHT)
    _txt(s1,team_name,Inches(1.1),Inches(2.05),Inches(8),Inches(.7),sz=24,italic=True,color=_rgb("#A9CCE3"))
    _txt(s1,f"{master.get('quarter','')}　{month_label}",Inches(1.1),Inches(2.78),Inches(7),Inches(.55),sz=18,color=_rgb("#A9CCE3"))
    _txt(s1,f"作成日：{datetime.date.today().strftime('%Y年%m月%d日')}",Inches(1.1),H*.65,Inches(7),Inches(.5),sz=14,color=C_SLAT)
    bx=Inches(1.1)
    for m in [p["member"] for p in all_plans if p.get("member")][:4]:
        pal=mpal(m); _rect(s1,bx,H*.77,Inches(2.55),Inches(.44),_rgb(pal["main"]))
        _txt(s1,m,bx+Cm(.1),H*.77+Cm(.05),Inches(2.35),Inches(.38),sz=13,bold=True,color=C_WHT,align=PP_ALIGN.CENTER)
        bx+=Inches(2.7)
    for rx,ry,rw in [(W-Inches(3.0),Inches(.1),Inches(2.8)),(W-Inches(5.3),H*.28,Inches(1.2))]:
        c=s1.shapes.add_shape(9,rx,ry,rw,rw); c.fill.solid(); c.fill.fore_color.rgb=C_GRN; c.line.fill.background()

    # OKRサマリー
    s2=prs.slides.add_slide(blank); _rect(s2,0,0,W,Inches(.95),C_DEEP)
    _txt(s2,"チーム OKR サマリー",Inches(.5),Inches(.1),Inches(11),Inches(.72),sz=26,bold=True,color=C_WHT)
    _txt(s2,f"🌟 {obj}",Inches(.5),Inches(1.1),W-Inches(1.0),Inches(.75),sz=15,bold=True,color=C_DEEP)
    ky=Inches(1.95); kh=Inches(.58)
    for i,kr in enumerate(krs[:3]):
        col=_rgb(KR_COLORS[i%len(KR_COLORS)]); _rect(s2,Inches(.5),ky,Inches(.06),kh,col)
        _rect(s2,Inches(.62),ky,W-Inches(1.12),kh,_rgb("#EBF5FB"))
        _txt(s2,f"{kr['label']}　{kr.get('text','')}",Inches(.75),ky+Cm(.12),W-Inches(1.4),kh-Cm(.2),sz=13,color=C_SLAT)
        ky+=kh+Inches(.12)
    card_w=(W-Inches(.8)-Inches(.25)*3)/4; cx=Inches(.4); cy=Inches(3.55); ch=H-cy-Inches(.2)
    for plan in all_plans[:4]:
        m=plan.get("member",""); pal=mpal(m)
        _rect(s2,cx,cy,card_w,ch,C_SNOW,border=_rgb("#D6DBDF")); _rect(s2,cx,cy,card_w,Inches(.42),_rgb(pal["main"]))
        _txt(s2,m,cx+Cm(.15),cy+Cm(.05),card_w-Cm(.3),Inches(.38),sz=11,bold=True,color=C_WHT,align=PP_ALIGN.CENTER)
        iy=cy+Inches(.5)
        for item in plan.get("items",[]):
            for wa in item.get("wall_actions",[]):
                for a in wa.get("actions",[])[:1]:
                    if a.get("text","").strip() and iy<cy+ch-Inches(.1):
                        _txt(s2,f"• {a['text'][:28]}",cx+Cm(.2),iy,card_w-Cm(.4),Inches(.45),sz=8,color=C_SLAT)
                        iy+=Inches(.38)
        cx+=card_w+Inches(.25)

    # メンバー別詳細
    for plan in all_plans:
        m=plan.get("member","不明"); pal=mpal(m)
        sm=prs.slides.add_slide(blank)
        _rect(sm,0,0,W,Inches(.95),_rgb(pal["main"])); _rect(sm,0,0,Inches(.05),H,_rgb(pal["main"]))
        _txt(sm,f"{m}　詳細アクション計画",Inches(.5),Inches(.1),Inches(10),Inches(.72),sz=24,bold=True,color=C_WHT)
        _txt(sm,f"🌟 {obj[:65]}",Inches(.5),Inches(1.05),W-Inches(1.0),Inches(.45),sz=11,italic=True,color=C_SLAT)
        hdrs=["KR","指標（Key Result）","課題","アクション","開始日","終了日"]
        col_ws=[Inches(.55),Inches(2.0),Inches(2.5),Inches(3.5),Inches(1.1),Inches(1.1)]
        row_h=Inches(.50); tbl_x=Inches(.4); tbl_y=Inches(1.6)
        cx=tbl_x
        for hdr,cw in zip(hdrs,col_ws):
            _rect(sm,cx,tbl_y,cw,row_h,_rgb(pal["main"]))
            _txt(sm,hdr,cx+Cm(.12),tbl_y+Cm(.08),cw-Cm(.24),row_h,sz=9,bold=True,color=C_WHT,align=PP_ALIGN.CENTER)
            cx+=cw
        row_bgs=[C_SNOW,RGBColor(0xFF,0xFF,0xFF)]; ri=0
        for item in plan.get("items",[]):
            for wa in item.get("wall_actions",[]):
                for a in [x for x in wa.get("actions",[]) if x.get("text","").strip()]:
                    cy_=tbl_y+row_h*(ri+1)
                    if cy_+row_h>H-Inches(.15): break
                    vals=[item.get("kr_label",""),item.get("kr_text","")[:35],wa.get("wall_text","")[:38],a.get("text","")[:50],a.get("start",""),a.get("end","")]
                    cx=tbl_x; bg=row_bgs[ri%2]
                    for ci,(val,cw) in enumerate(zip(vals,col_ws)):
                        _rect(sm,cx,cy_,cw,row_h,bg,border=_rgb("#D5D8DC"))
                        _txt(sm,val,cx+Cm(.15),cy_+Cm(.07),cw-Cm(.3),row_h,sz=9,
                             color=_rgb(pal["main"]) if ci==0 else C_SLAT,
                             align=PP_ALIGN.CENTER if ci in(0,4,5) else PP_ALIGN.LEFT)
                        cx+=cw
                    ri+=1

    # ガントチャート
    sg=prs.slides.add_slide(blank); _rect(sg,0,0,W,Inches(.95),C_NAVY)
    _txt(sg,"統合ガントチャート",Inches(.5),Inches(.1),Inches(11),Inches(.72),sz=26,bold=True,color=C_WHT)
    if gantt_png:
        sg.shapes.add_picture(io.BytesIO(gantt_png),Inches(.3),Inches(1.1),width=W-Inches(.6),height=H-Inches(1.3))
    else:
        _txt(sg,"アクションデータがありません",Inches(1.0),Inches(3.0),W-Inches(2.0),Inches(.8),sz=18,color=C_GRAY,align=PP_ALIGN.CENTER)

    for sl in prs.slides:
        _txt(sl,f"{team_name}　{month_label}",Inches(.4),H-Inches(.32),W-Inches(.8),Inches(.28),sz=8,italic=True,color=C_GRAY,align=PP_ALIGN.RIGHT)

    buf=io.BytesIO(); prs.save(buf); return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# メイン
# ══════════════════════════════════════════════════════════════════════════════

def main():
    _init_session()
    team_name  = CFG["team_name"]
    month_str  = st.session_state.month_str
    month_disp = datetime.date.fromisoformat(month_str+"-01").strftime("%Y年%m月")

    # マスターデータをセッションキャッシュ（初回のみSheets読み込み）
    if "cached_master" not in st.session_state:
        st.session_state.cached_master = io_get_master()
    master = st.session_state.cached_master

    with st.sidebar:
        st.markdown(f"## 🌟 OKR管理\n**{team_name}**")
        st.markdown("---")

        # 接続状況（初回のみ確認）
        if "sheets_status" not in st.session_state:
            _c = get_gsheet_client()
            if isinstance(_c, str) and _c.startswith("ERROR:"):
                st.session_state.sheets_status = ("error", _c[6:120])
            elif _c and CFG.get("sheet_id"):
                st.session_state.sheets_status = ("ok", "")
            else:
                st.session_state.sheets_status = ("local", "")
        _status, _msg = st.session_state.sheets_status
        if _status == "error":
            st.markdown('<div class="local-badge" style="background:#FDEDEC;border-color:#F1948A;color:#922B21;">❌ Sheets接続失敗</div>', unsafe_allow_html=True)
            if _msg: st.caption(_msg)
        elif _status == "ok":
            st.markdown('<div class="local-badge" style="background:#EAF7EE;border-color:#82E0AA;color:#196F3D;">✅ Sheets接続済み</div>', unsafe_allow_html=True)
        else:
            st.markdown('<div class="local-badge">⚙️ ローカル保存モード</div>', unsafe_allow_html=True)

        st.markdown(f"### 📅 対象月\n**{month_disp}**")
        st.caption(f"四半期：{master.get('quarter','未設定')}")
        st.markdown("---")
        st.markdown("### メンバー")
        for i, m in enumerate(MEMBERS):
            c = PALETTE[i % len(PALETTE)]["main"]
            st.markdown(
                f'<div style="display:flex;align-items:center;gap:7px;margin:3px 0;">'
                f'<div style="width:10px;height:10px;border-radius:50%;background:{c};flex-shrink:0;"></div>'
                f'<span style="font-size:.82rem;color:#C8E0FF;">{m}</span></div>',
                unsafe_allow_html=True,
            )
        if st.session_state.admin_auth:
            st.markdown("---")
            st.markdown('<span style="color:#F9E08B;font-size:.78rem;">🔐 管理者モード ON</span>', unsafe_allow_html=True)
            if st.button("ログアウト", use_container_width=True):
                st.session_state.admin_auth = False; st.rerun()

    render_north_star(master)

    tab_home, tab_strategy, tab_plan, tab_dash, tab_task = st.tabs([
        "🏠 HOME", "🏛️ STRATEGY",
        "📝 PLAN", "📊 DASHBOARD",
        "🎫 タスク起票",
    ])

    with tab_home:
        render_home()
    with tab_strategy:
        st.subheader("🏛️ 四半期 OKR 作成ワークショップ")
        render_strategy(master)
    with tab_plan:
        st.subheader(f"📝 {month_disp} チームプラン")
        render_plan(master)
    with tab_dash:
        st.subheader(f"📊 {month_disp}　チーム統合ダッシュボード")
        render_dashboard(master)
    with tab_task:
        st.subheader("🎫 タスク起票")
        render_task_ticket(master, month_str)


if __name__ == "__main__":
    main()
